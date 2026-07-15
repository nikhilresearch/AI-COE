"""
BIMTECH AI Build Bootcamp & Hackathon 2026 – PowerPoint Generator
Run: python create_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colours ──────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x06, 0x0d, 0x2b)
BLUE_BRIGHT = RGBColor(0x1a, 0x4f, 0xff)
CYAN        = RGBColor(0x00, 0xd4, 0xff)
PURPLE      = RGBColor(0x7c, 0x3a, 0xed)
WHITE       = RGBColor(0xff, 0xff, 0xff)
GOLD        = RGBColor(0xf5, 0xc8, 0x42)
LIGHT_BLUE  = RGBColor(0xc4, 0xaa, 0xff)
GREEN       = RGBColor(0x4a, 0xde, 0x80)
GREY_LIGHT  = RGBColor(0xcc, 0xdd, 0xee)

# Slide dimensions (widescreen 16:9)
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helper Functions ──────────────────────────────────────────────────────────

def blank_slide(prs):
    """Return a slide with completely blank layout."""
    layout = prs.slide_layouts[6]   # Blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color: RGBColor):
    """Solid-colour slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width_pt=0.75,
             transparency=0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = Pt(line_width_pt) if line_color else Pt(0)
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb

def add_gradient_rect(slide, left, top, width, height):
    """Approx gradient using two overlapping rects."""
    r1 = add_rect(slide, left, top, width, height, BLUE_BRIGHT)
    r2 = add_rect(slide, left, top, Inches(2), height, PURPLE)
    return r1

def multi_para_textbox(slide, left, top, width, height, paras,
                       font_name="Calibri"):
    """
    paras = list of dicts:
      text, font_size, bold, color, align, italic, space_before
    """
    from pptx.util import Pt
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for p in paras:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        if p.get('space_before'):
            para.space_before = Pt(p['space_before'])
        run = para.add_run()
        run.text = p.get('text', '')
        run.font.size = Pt(p.get('font_size', 16))
        run.font.bold = p.get('bold', False)
        run.font.italic = p.get('italic', False)
        run.font.color.rgb = p.get('color', WHITE)
        run.font.name = font_name
    return tb

# ── Slide 1 ───────────────────────────────────────────────────────────────────
slide1 = blank_slide(prs)
fill_bg(slide1, NAVY)

# Left accent bar
add_rect(slide1, Inches(0), Inches(0), Inches(.18), H, BLUE_BRIGHT)

# Top glow strip
add_rect(slide1, Inches(.18), Inches(0), W, Inches(.06), CYAN)

# Collaboration badge background
badge = add_rect(slide1, Inches(.5), Inches(.4), Inches(3.5), Inches(.45),
                 RGBColor(0x00, 0x22, 0x44))
add_textbox(slide1, "✦  BIMTECH × Reskilll  ✦",
            Inches(.5), Inches(.4), Inches(3.5), Inches(.45),
            font_size=10, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# Title
multi_para_textbox(slide1, Inches(.5), Inches(1.0), Inches(8), Inches(2.2), [
    {"text": "AI Build Bootcamp", "font_size": 46, "bold": True, "color": WHITE},
    {"text": "& Hackathon 2026",  "font_size": 46, "bold": True, "color": CYAN},
])

# Subtitle
add_textbox(slide1,
    "Empowering Future Managers with AI Skills",
    Inches(.5), Inches(3.2), Inches(9), Inches(.55),
    font_size=19, bold=False, color=GREY_LIGHT)

# Date pill
add_rect(slide1, Inches(.5), Inches(3.9), Inches(5.2), Inches(.55),
         RGBColor(0x1a, 0x08, 0x50), PURPLE, 1.0)
add_textbox(slide1,
    "📅  7–8 August 2026   |   9:00 AM – 7:00 PM",
    Inches(.5), Inches(3.9), Inches(5.2), Inches(.55),
    font_size=13, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Divider
add_rect(slide1, Inches(.5), Inches(4.65), Inches(1.2), Inches(.05), CYAN)

# Two info cards
for i, (title, body) in enumerate([
    ("🤔  Why This Initiative?",
     "AI is transforming every business function — Finance, Marketing, HR,\n"
     "Operations & Strategy. Tomorrow's managers will not only use AI but\n"
     "also build AI-enabled business solutions."),
    ("🎯  Objective",
     "Equip every first-year MBA student with hands-on experience in\n"
     "designing and building AI-powered business applications."),
]):
    left = Inches(.5 + i * 5.9)
    add_rect(slide1, left, Inches(4.85), Inches(5.6), Inches(2.1),
             RGBColor(0x0a, 0x17, 0x45), CYAN, 0.75)
    add_textbox(slide1, title, left + Inches(.18), Inches(4.95),
                Inches(5.2), Inches(.4),
                font_size=11, bold=True, color=CYAN)
    add_textbox(slide1, body, left + Inches(.18), Inches(5.42),
                Inches(5.2), Inches(1.4),
                font_size=10.5, color=GREY_LIGHT)

# Slide number
add_textbox(slide1, "01", Inches(12.8), Inches(.1), Inches(.5), Inches(.4),
            font_size=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# ── Slide 2 ───────────────────────────────────────────────────────────────────
slide2 = blank_slide(prs)
fill_bg(slide2, NAVY)
add_rect(slide2, Inches(0), Inches(0), Inches(.18), H, PURPLE)
add_rect(slide2, Inches(.18), Inches(0), W, Inches(.06), CYAN)

add_textbox(slide2, "✦  STUDENT IMPACT  ✦",
            Inches(.4), Inches(.3), Inches(3.5), Inches(.38),
            font_size=9, bold=True, color=CYAN)

multi_para_textbox(slide2, Inches(.4), Inches(.8), Inches(8), Inches(1.6), [
    {"text": "Why This Matters", "font_size": 36, "bold": True, "color": WHITE},
    {"text": "to Our Students",   "font_size": 36, "bold": True, "color": CYAN},
])

add_rect(slide2, Inches(.4), Inches(1.9), Inches(1.0), Inches(.06), CYAN)

# Left card – Students will learn
add_rect(slide2, Inches(.4), Inches(2.1), Inches(5.8), Inches(4.2),
         RGBColor(0x08, 0x14, 0x40), CYAN, 0.75)
add_textbox(slide2, "🧠  Students Will Learn To",
            Inches(.55), Inches(2.2), Inches(5.4), Inches(.4),
            font_size=11, bold=True, color=CYAN)

bullets_learn = [
    "→  Convert business problems into AI solutions",
    "→  Work collaboratively in multidisciplinary teams",
    "→  Build AI-enabled prototypes without extensive coding",
    "→  Develop innovation & design-thinking skills",
    "→  Present solutions before industry experts",
]
for idx, b in enumerate(bullets_learn):
    add_textbox(slide2, b,
                Inches(.6), Inches(2.72 + idx * .57), Inches(5.4), Inches(.5),
                font_size=11, color=GREY_LIGHT)

# Right card – Outcomes
add_rect(slide2, Inches(6.5), Inches(2.1), Inches(6.4), Inches(4.2),
         RGBColor(0x08, 0x14, 0x40), PURPLE, 0.75)
add_textbox(slide2, "🏆  Outcomes",
            Inches(6.65), Inches(2.2), Inches(5.9), Inches(.4),
            font_size=11, bold=True, color=CYAN)

outcomes = [
    "✔  Hands-on AI Exposure",
    "✔  Practical Problem-Solving",
    "✔  Industry Interaction",
    "✔  Innovation Mindset",
    "✔  Portfolio-Worthy Projects",
]
for idx, o in enumerate(outcomes):
    add_textbox(slide2, o,
                Inches(6.65), Inches(2.72 + idx * .57), Inches(5.5), Inches(.5),
                font_size=11, color=GREEN)

# Quote band
add_rect(slide2, Inches(.4), Inches(6.5), Inches(12.5), Inches(.68),
         RGBColor(0x0a, 0x1a, 0x5c), CYAN, 0.75)
add_textbox(slide2,
    '"Learn AI by Building, Not Just by Reading."',
    Inches(.55), Inches(6.55), Inches(12.2), Inches(.6),
    font_size=14, bold=True, italic=True, color=CYAN, align=PP_ALIGN.CENTER)

add_textbox(slide2, "02", Inches(12.8), Inches(.1), Inches(.5), Inches(.4),
            font_size=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# ── Slide 3 ───────────────────────────────────────────────────────────────────
slide3 = blank_slide(prs)
fill_bg(slide3, NAVY)
add_rect(slide3, Inches(0), Inches(0), Inches(.18), H, BLUE_BRIGHT)
add_rect(slide3, Inches(.18), Inches(0), W, Inches(.06), CYAN)

add_textbox(slide3, "✦  EVENT STRUCTURE  ✦",
            Inches(.4), Inches(.3), Inches(3.5), Inches(.38),
            font_size=9, bold=True, color=CYAN)

multi_para_textbox(slide3, Inches(.4), Inches(.75), Inches(10), Inches(1.2), [
    {"text": "AI Build Bootcamp  ", "font_size": 34, "bold": True, "color": WHITE},
    {"text": "& Hackathon",         "font_size": 34, "bold": True, "color": CYAN},
])

add_rect(slide3, Inches(.4), Inches(1.75), Inches(1.0), Inches(.06), CYAN)

# 3 day columns
col_configs = [
    {
        "header_color": BLUE_BRIGHT,
        "title_label":  "📅 Day 1 — 7 August",
        "title":        "AI Bootcamp",
        "bullets": [
            "Introduction to AI tools",
            "Guided hands-on sessions",
            "Team formation",
            "Problem statement briefing",
        ],
        "bullet_color": GREY_LIGHT,
    },
    {
        "header_color": PURPLE,
        "title_label":  "🚀 Day 2 — 8 August",
        "title":        "AI Hackathon",
        "bullets": [
            "Prototype development",
            "Mentoring support",
            "Project submission",
            "Evaluation by Reskilll",
        ],
        "bullet_color": GREY_LIGHT,
    },
    {
        "header_color": RGBColor(0xb4, 0x5a, 0x00),
        "title_label":  "🏅 19 August — Pitch Day",
        "title":        "Top 10 Teams",
        "bullets": [
            "Present before Industry Experts",
            "Advisory Committee",
            "Faculty Jury",
            "⭐ Top 3 Teams Rewarded",
        ],
        "bullet_color": GOLD,
    },
]

col_left = [Inches(.4), Inches(4.65), Inches(8.9)]
col_w = Inches(4.0)

for i, cfg in enumerate(col_configs):
    cl = col_left[i]
    # card bg
    add_rect(slide3, cl, Inches(2.0), col_w, Inches(4.9),
             RGBColor(0x08, 0x12, 0x3e), cfg['header_color'], 0.75)
    # label
    add_textbox(slide3, cfg['title_label'], cl + Inches(.15), Inches(2.1),
                col_w - Inches(.3), Inches(.38),
                font_size=9.5, bold=True, color=CYAN)
    # title
    add_textbox(slide3, cfg['title'], cl + Inches(.15), Inches(2.55),
                col_w - Inches(.3), Inches(.5),
                font_size=16, bold=True, color=WHITE)
    # bullets
    for bi, b in enumerate(cfg['bullets']):
        add_textbox(slide3, "▸  " + b,
                    cl + Inches(.2), Inches(3.15 + bi * .65),
                    col_w - Inches(.35), Inches(.58),
                    font_size=11, color=cfg['bullet_color'])

add_textbox(slide3, "03", Inches(12.8), Inches(.1), Inches(.5), Inches(.4),
            font_size=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# ── Slide 4 ───────────────────────────────────────────────────────────────────
slide4 = blank_slide(prs)
fill_bg(slide4, NAVY)
add_rect(slide4, Inches(0), Inches(0), Inches(.18), H, PURPLE)
add_rect(slide4, Inches(.18), Inches(0), W, Inches(.06), CYAN)

add_textbox(slide4, "✦  YOUR ROLE MATTERS  ✦",
            Inches(.4), Inches(.3), Inches(4), Inches(.38),
            font_size=9, bold=True, color=CYAN)

add_textbox(slide4, "Role of Programme Chairs",
            Inches(.4), Inches(.78), Inches(9), Inches(.72),
            font_size=34, bold=True, color=WHITE)
add_textbox(slide4,
    "Programme Chairs play a critical role in making this initiative successful.",
    Inches(.4), Inches(1.6), Inches(12.5), Inches(.44),
    font_size=12, color=GREY_LIGHT)

add_rect(slide4, Inches(.4), Inches(2.12), Inches(1.0), Inches(.06), CYAN)

# Left – Support Required
add_rect(slide4, Inches(.4), Inches(2.3), Inches(5.8), Inches(4.7),
         RGBColor(0x08, 0x14, 0x40), BLUE_BRIGHT, 0.75)
add_textbox(slide4, "🤝  Support Required",
            Inches(.55), Inches(2.42), Inches(5.4), Inches(.4),
            font_size=11, bold=True, color=CYAN)

support = [
    "Encourage maximum student participation",
    "Ensure students attend as per batch allocation",
    "Minimise class disruptions during the event",
    "Motivate students to treat it as a learning opportunity",
    "Coordinate with faculty mentors wherever required",
]
for idx, s in enumerate(support):
    add_textbox(slide4, "→  " + s,
                Inches(.6), Inches(2.95 + idx * .62), Inches(5.4), Inches(.55),
                font_size=10.5, color=GREY_LIGHT)

# Right – Student Benefits
add_rect(slide4, Inches(6.5), Inches(2.3), Inches(6.4), Inches(4.7),
         RGBColor(0x15, 0x08, 0x40), PURPLE, 0.75)
add_textbox(slide4, "🎁  Student Benefits",
            Inches(6.65), Inches(2.42), Inches(6.0), Inches(.4),
            font_size=11, bold=True, color=GOLD)

benefits = [
    ("🏆", "Joint BIMTECH–Reskilll Certificate"),
    ("🏭", "Industry Exposure"),
    ("💼", "AI Portfolio Project"),
    ("🥇", "Recognition for Top Teams"),
    ("🎙️", "Present Before Industry Leaders"),
]
for idx, (icon, label) in enumerate(benefits):
    add_rect(slide4, Inches(6.65), Inches(2.98 + idx * .72), Inches(5.9), Inches(.56),
             RGBColor(0x20, 0x10, 0x50), GOLD, 0.5)
    add_textbox(slide4, icon + "  " + label,
                Inches(6.80), Inches(3.0 + idx * .72), Inches(5.7), Inches(.52),
                font_size=11, color=GOLD)

add_textbox(slide4, "04", Inches(12.8), Inches(.1), Inches(.5), Inches(.4),
            font_size=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# ── Slide 5 ───────────────────────────────────────────────────────────────────
slide5 = blank_slide(prs)
fill_bg(slide5, NAVY)
add_rect(slide5, Inches(0), Inches(0), Inches(.18), H, CYAN)
add_rect(slide5, Inches(.18), Inches(0), W, Inches(.06), PURPLE)

add_textbox(slide5, "✦  BIGGER PICTURE  ✦",
            Inches(.4), Inches(.3), Inches(3.5), Inches(.38),
            font_size=9, bold=True, color=CYAN)

multi_para_textbox(slide5, Inches(.4), Inches(.78), Inches(9), Inches(1.4), [
    {"text": "Building the AI Culture", "font_size": 34, "bold": True, "color": WHITE},
    {"text": "at BIMTECH",              "font_size": 34, "bold": True, "color": CYAN},
])

# Quote card
add_rect(slide5, Inches(.4), Inches(2.0), Inches(12.5), Inches(.9),
         RGBColor(0x0a, 0x1a, 0x5c), CYAN, 0.75)
add_textbox(slide5,
    '"This is more than a Hackathon. It is the first flagship initiative of the AI Centre of Excellence."',
    Inches(.6), Inches(2.08), Inches(12.1), Inches(.72),
    font_size=13, bold=True, italic=True, color=CYAN, align=PP_ALIGN.CENTER)

# Vision card
add_rect(slide5, Inches(.4), Inches(3.1), Inches(5.8), Inches(1.3),
         RGBColor(0x08, 0x14, 0x40), PURPLE, 0.75)
add_textbox(slide5, "🌟  Vision",
            Inches(.55), Inches(3.2), Inches(5.4), Inches(.38),
            font_size=11, bold=True, color=GOLD)
add_textbox(slide5,
    "Create AI-literate management professionals capable of solving\nreal business challenges using AI.",
    Inches(.55), Inches(3.65), Inches(5.4), Inches(.7),
    font_size=11, color=GREY_LIGHT)

# Journey flow
add_rect(slide5, Inches(.4), Inches(4.65), Inches(12.5), Inches(1.05),
         RGBColor(0x06, 0x12, 0x38), BLUE_BRIGHT, 0.75)
add_textbox(slide5, "🗺️  Our Journey",
            Inches(.55), Inches(4.72), Inches(3), Inches(.38),
            font_size=11, bold=True, color=CYAN)

journey_steps = [
    "AI Bootcamp", "➜", "AI Hackathon", "➜", "Pitch Day",
    "➜", "AI Centre of Excellence", "➜", "Industry Projects", "➜", "🚀 AI Innovation Ecosystem"
]
step_left = Inches(.5)
for step in journey_steps:
    is_arrow = step == "➜"
    color = GREY_LIGHT if is_arrow else (GOLD if "🚀" in step else CYAN)
    width = Inches(.35) if is_arrow else Inches(1.55)
    add_textbox(slide5, step,
                step_left, Inches(5.2), width, Inches(.42),
                font_size=10, bold=(not is_arrow), color=color,
                align=PP_ALIGN.CENTER)
    step_left += width + Inches(.04)

add_textbox(slide5, "05", Inches(12.8), Inches(.1), Inches(.5), Inches(.4),
            font_size=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__),
                        "BIMTECH_AI_Hackathon_2026.pptx")
prs.save(out_path)
print("Saved: " + out_path)
